import streamlit as st
import openai
import docx
import pandas as pd
import PyPDF2
import pptx
from io import BytesIO
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

# OpenAI API Configuration
openai.api_key = "14560021aaf84772835d76246b53397a"
openai.api_base = "https://amrxgenai.openai.azure.com/"
openai.api_type = 'azure'
openai.api_version = '2024-02-15-preview'
deployment_name = 'gpt'

# Session state init
if "messages" not in st.session_state:
    st.session_state.messages = []
if "uploaded_file" not in st.session_state:
    st.session_state.uploaded_file = None
if "extracted_text" not in st.session_state:
    st.session_state.extracted_text = ""
if "new_upload" not in st.session_state:
    st.session_state.new_upload = False
if "sidebar_query" not in st.session_state:
    st.session_state.sidebar_query = ""

# Extract text function
def extract_text(file):
    text = ""
    if file.type == "application/pdf":
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file.type == "text/csv":
        df = pd.read_csv(file)
        text += df.to_string()
    elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(file)
        text += df.to_string()
    return text.strip()

# Create DOCX history
def create_docx():
    doc = Document()
    doc.add_heading("Chatbot Conversation History", level=1).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    for msg in st.session_state.messages:
        p = doc.add_paragraph()
        run = p.add_run(msg["role"].capitalize() + "\n")
        run.bold = True
        run.font.size = Pt(14)
        if "file" in msg and msg["file"]:
            doc.add_paragraph(f"📄 {msg['file']}")
        p.add_run(msg["content"]).font.size = Pt(12)
        p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        doc.add_paragraph("\n")
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# SIDEBAR
with st.sidebar:
    st.image("image.png", width=150)
    st.markdown(
        f'<a href="https://fanalysis-uhevotzeiesw3nczetzsmy.streamlit.app/" target="_blank">'
        f'<img src="image.png" width="150"></a>',
        unsafe_allow_html=True
    )
    st.image("https://www.innominds.com/hubfs/Innominds-201612/img/nav/Innominds-Logo.png", width=200)
    st.title("How Can I Assist You? 🤖")

    # Upload button logic
    if st.button("➕ Upload File"):
        st.session_state.new_upload = True

    if st.session_state.new_upload:
        uploaded_file = st.file_uploader("Upload File", type=["pdf", "docx", "pptx", "csv", "xlsx"], key="file_uploader")
        if uploaded_file:
            st.session_state.uploaded_file = uploaded_file.name
            st.session_state.extracted_text = extract_text(uploaded_file)
            st.session_state.messages.insert(0, {
                "role": "system",
                "content": f"Document Context:\n\n{st.session_state.extracted_text}"
            })
            st.session_state.new_upload = False

    if st.session_state.uploaded_file:
        col1, col2 = st.columns([5, 1])
        with col1:
            st.markdown(f"📂 Uploaded: **{st.session_state.uploaded_file}**")
        with col2:
            if st.button("❌ Clear File"):
                st.session_state.uploaded_file = None
                st.session_state.extracted_text = ""

    # Query bar in sidebar
    query = st.text_input("🔍 Ask something...", key="sidebar_query")
    if query:
        combined_prompt = query
        if st.session_state.uploaded_file:
            combined_prompt = f"Here is a document that provides context:\n\n{st.session_state.extracted_text}\n\nNow, based on this document, answer the following:\n{query}"

        st.session_state.messages.append({
            "role": "user",
            "content": query,
            "file": st.session_state.uploaded_file if st.session_state.uploaded_file else None
        })

        response = openai.ChatCompletion.create(
            engine=deployment_name,
            messages=st.session_state.messages,
            temperature=0,
            max_tokens=2000
        )
        ai_response = response["choices"][0]["message"]["content"]

        st.session_state.messages.append({
            "role": "assistant",
            "content": ai_response
        })

        st.experimental_rerun()

    # Download chat history
    if st.button("Download Chat History"):
        docx_file = create_docx()
        st.download_button(
            label="Download Chat History",
            data=docx_file,
            file_name="chat_history.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

# MAIN PAGE – chat messages
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        if "file" in message and message["file"]:
            st.markdown(f"📄 **{message['file']}**")
        st.markdown(message["content"])

# Optional: chat input at bottom
user_input = st.chat_input("Enter your query:")
if user_input and user_input.strip():
    combined_prompt = user_input
    if st.session_state.uploaded_file:
        combined_prompt = f"Here is a document that provides context:\n\n{st.session_state.extracted_text}\n\nNow, based on this document, answer the following:\n{user_input}"

    st.session_state.messages.append({
        "role": "user",
        "content": user_input,
        "file": st.session_state.uploaded_file if st.session_state.uploaded_file else None
    })

    response = openai.ChatCompletion.create(
        engine=deployment_name,
        messages=st.session_state.messages,
        temperature=0,
        max_tokens=2000
    )
    ai_response = response["choices"][0]["message"]["content"]

    st.session_state.messages.append({
        "role": "assistant",
        "content": ai_response
    })

    st.rerun()
