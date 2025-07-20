import streamlit as st
import openai
import docx
import pandas as pd
import PyPDF2
import pptx
from io import BytesIO
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT


st.set_page_config(
    page_title="Innominds GEN-AI",  # must be non-empty!
    page_icon="🤖",
    layout="wide"
)



# --- PCB/BCB Prompt Template ---
PCB_BCB_TEMPLATE ="""
If the question is about providing / generating the requirements of a bonus, Please use the following section headers and content.
	Introduction
	- provide the bonus overview in this section
	Commission configurations
	- use the records with the type as "commissions configurations".
	- Show the parameter configuration and description as bullet points
	Rate card configurations
	- use the records with the type as "Rate card configurations".
	- Show the parameter configuration and description as bullet points
	Process
	- use the records with the type as "Process".
	- use the Key parameter definition text to describe the process
 
If the question is about Comparison of bonuses, show the output in Tabular format.
 
If the question is about generating requirements of a specific bonus or all the bonuses for a country, based on the parameters of existing countries, do the following:
- Convert the currency to the country for which the bonus requirements are being created
- Generate the bonus requirement as per the existing countries bonus structures.
"""

# --- Sidebar content ---
with st.sidebar:
    st.title("How Can I Assist You?🤖")
    user_input = st.chat_input("Enter your query:")

    if st.button("➕ Upload File"):
        st.session_state.new_upload = True

    if st.session_state.get("new_upload", False):
        uploaded_file = st.file_uploader("Upload File", type=["pdf", "docx", "pptx", "csv", "xlsx"], key="file_uploader")

        if uploaded_file:
            st.session_state.uploaded_file = uploaded_file.name
            st.session_state.extracted_text = ""

            def extract_text(file):
                text = ""
                if file.type == "application/pdf":
                    reader = PyPDF2.PdfReader(file)
                    for page in reader.pages:
                        extracted = page.extract_text()
                        if extracted:
                            text += extracted + "\n"
                elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    doc = Document(file)
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    ppt = pptx.Presentation(file)
                    for slide in ppt.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                elif file.type == "text/csv":
                    df = pd.read_csv(file)
                    text += df.to_string()
                elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                    df = pd.read_excel(file)
                    text += df.to_string()
                return text.strip()

            st.session_state.extracted_text = extract_text(uploaded_file)
            st.session_state.new_upload = False

    if st.session_state.get("uploaded_file"):
        st.markdown(f"📂 **Uploaded File:** {st.session_state.uploaded_file}")
        if st.button("❌ Clear Uploaded File", use_container_width=True):
            st.session_state.uploaded_file = None
            st.session_state.extracted_text = ""

# --- Main content ---
st.image("https://www.innominds.com/hubfs/Innominds-201612/img/nav/Innominds-Logo.png", width=200)

# Initialize session state
if "messages" not in st.session_state:
    st.session_state.messages = []
if "uploaded_file" not in st.session_state:
    st.session_state.uploaded_file = None
if "extracted_text" not in st.session_state:
    st.session_state.extracted_text = ""
if "new_upload" not in st.session_state:
    st.session_state.new_upload = False

# Display chat messages
for message in st.session_state.messages:
    if message["role"] != "system":
        with st.chat_message(message["role"]):
            if "file" in message and message["file"]:
                st.markdown(f"📄 **{message['file']}**")
            st.markdown(message["content"])

# OpenAI API Configuration
openai.api_key = "14560021aaf84772835d76246b53397a"
openai.api_base = "https://amrxgenai.openai.azure.com/"
openai.api_type = 'azure'
openai.api_version = '2024-02-15-preview'
deployment_name = 'gpt'

# Chat handling
if user_input and user_input.strip():
    chat_messages = st.session_state.messages.copy()

    # Determine if PCB or BCB is mentioned
    pcb_bcb_in_query = "pcb" in user_input.lower() or "bcb" in user_input.lower()

    # Compose the prompt
    if st.session_state.extracted_text:
        if not any(m["role"] == "system" for m in chat_messages):
            system_prompt = f"Here is a document that provides context:\n\n{st.session_state.extracted_text}\n\nNow, based on this document, answer the following:\n{user_input}"
            if pcb_bcb_in_query:
                system_prompt += PCB_BCB_TEMPLATE + "\n\n"
            system_prompt += f"Now, based on this document, answer the following:\n{user_input}"
            chat_messages.insert(0, {
                "role": "system",
                "content": system_prompt
            })

    chat_messages.append({
        "role": "user",
        "content": user_input
    })

    response = openai.ChatCompletion.create(
        engine=deployment_name,
        messages=chat_messages,
        temperature=0,
        max_tokens=4095
    )

    ai_response = response["choices"][0]["message"]["content"]

    st.session_state.messages.append({
        "role": "user",
        "content": user_input,
        "file": st.session_state.uploaded_file
    })

    st.session_state.messages.append({
        "role": "assistant",
        "content": ai_response
    })

    st.rerun()

# --- Download DOCX button ---
def create_docx():
    doc = Document()
    doc.add_heading("Chatbot Conversation History", level=1).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    for msg in st.session_state.messages:
        if msg["role"] == "system":
            continue
        p = doc.add_paragraph()
        run = p.add_run(msg["role"].capitalize() + "\n")
        run.bold = True
        run.font.size = Pt(14)
        if "file" in msg and msg["file"]:
            doc.add_paragraph(f"📄 {msg['file']}")
        p.add_run(msg["content"]).font.size = Pt(12)
        p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        doc.add_paragraph("\n")

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

if st.button("Download Chat History"):
    docx_file = create_docx()
    st.download_button(label="Download Chat History", data=docx_file, file_name="chat_history.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
