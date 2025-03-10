import streamlit as st
import openai
import docx
import pandas as pd
import PyPDF2
import pptx
from io import BytesIO
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import re


st.image("https://www.innominds.com/hubfs/Innominds-201612/img/nav/Innominds-Logo.png", width=200)

# OpenAI API Configuration
openai.api_key = "14560021aaf84772835d76246b53397a"
openai.api_base = "https://amrxgenai.openai.azure.com/"
openai.api_type = 'azure'
openai.api_version = '2024-02-15-preview'
deployment_name = 'gpt'

st.title("RESPONSES🤖")

# Initialize session state
if "messages" not in st.session_state:
    st.session_state.messages = []
if "uploaded_file" not in st.session_state:
    st.session_state.uploaded_file = None
if "extracted_text" not in st.session_state:
    st.session_state.extracted_text = ""
if "new_upload" not in st.session_state:
    st.session_state.new_upload = False

# Function to extract text from uploaded files (store in session but don't display)
def extract_text(file):
    text = ""
    if file.type == "application/pdf":
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file.type == "text/csv":
        df = pd.read_csv(file)
        text += df.to_string()
    elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(file)
        text += df.to_string()
    return text.strip()

# Display chat messages in a conversational format
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        if "file" in message and message["file"]:
            st.markdown(f"📄 **{message['file']}**")
        st.markdown(message["content"])

# Chat input section with "+" button for file upload
col1, col2 = st.columns([5, 1])
with col1:
    user_input = st.chat_input("Enter your query:")
with col2:
    if st.button("➕"):  # Clicking this opens the file uploader
        st.session_state.new_upload = True

# Persistent file uploader
if st.session_state.new_upload:
    uploaded_file = st.file_uploader("Upload File", type=["pdf", "docx", "pptx", "csv", "xlsx"], key="file_uploader")
    
    if uploaded_file:
        st.session_state.uploaded_file = uploaded_file.name
        st.session_state.extracted_text = extract_text(uploaded_file)  # Store extracted text but don't display it
        st.session_state.new_upload = False  # Reset after successful upload

# Display uploaded file info and Clear File button
if st.session_state.uploaded_file:
    col1, col2 = st.columns([5, 1])
    with col1:
        st.markdown(f"📂 Uploaded: **{st.session_state.uploaded_file}**")
    with col2:
        if st.button("❌ Clear File"):
            st.session_state.uploaded_file = None
            st.session_state.extracted_text = ""

if user_input and user_input.strip():
    combined_prompt = user_input
    # Include extracted text from uploaded file
    if st.session_state.extracted_text:
        combined_prompt = f"Context from uploaded file:\n{st.session_state.extracted_text}\n\nUser Query:\n{user_input}"
    if st.session_state.messages:  
        history = "\n\n".join([msg["content"] for msg in st.session_state.messages if msg["role"] == "assistant"])
        combined_prompt = f"Previous responses:\n{history}\n\nNow, based on this history, answer the following:\n{user_input}"
    
    # Append user message to history
    st.session_state.messages.append({"role": "user", "content": user_input, "file": st.session_state.uploaded_file if st.session_state.uploaded_file else None})
    
    # Generate response using OpenAI API
    response = openai.ChatCompletion.create(
        engine=deployment_name,
        messages=[{"role": "system", "content": combined_prompt}],
        temperature=0.7,
        max_tokens=2000
    )
    
    ai_response = re.sub(r'<br\s*/?>', '\n', response["choices"][0]["message"]["content"])  # Remove unwanted <br>
    
    # Append AI response to history
    st.session_state.messages.append({"role": "assistant", "content": ai_response})
    
    st.rerun()

# Function to create a DOCX file from chat history
def create_docx():
    doc = Document()
    doc.add_heading("Chatbot Conversation History", level=1).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    
    for msg in st.session_state.messages:
        p = doc.add_paragraph()
        run = p.add_run(msg["role"].capitalize() + "\n")
        run.bold = True
        run.font.size = Pt(14)
        if "file" in msg and msg["file"]:
            doc.add_paragraph(f"📄 {msg['file']}")
        p.add_run(msg["content"]).font.size = Pt(12)
        p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        doc.add_paragraph("\n")
    
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Download chat history as DOCX
if st.button("Download Chat History"):
    docx_file = create_docx()
    st.download_button(label="Download Chat History", data=docx_file, file_name="chat_history.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
