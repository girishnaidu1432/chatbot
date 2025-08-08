import streamlit as st
import openai
import docx
import pandas as pd
import PyPDF2
import pptx
from io import BytesIO
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

st.set_page_config(
    page_title="Innominds GEN-AI",
    page_icon="🤖",
    layout="wide"
)

# ---------------------------
# Remove padding & place logo at top
# ---------------------------
# Remove default padding so content starts at the very top
st.markdown(
    """
    <style>
        .block-container {
            padding-top: 0rem;
        }
    </style>
    """,
    unsafe_allow_html=True
)

# Full visible logo at the top
st.markdown(
    """
    <div style="display: flex; justify-content: flex-start; align-items: center; padding-top: 5px; padding-bottom: 10px;">
        <img src="https://www.innominds.com/hubfs/Innominds-201612/img/nav/Innominds-Logo.png" width="200">
    </div>
    """,
    unsafe_allow_html=True
)

# ---------------------------
# Preset prompts
# ---------------------------
PRESET_PROMPTS = [
    "What are the bonuses having the conditions for PANF, TANF or UNF region wise?",
    "Give the comparison matrix for SCNB and SCNB2 for poland",
    "give the comparison matrix for DVB1 for poland and spain",
    "Can you provide the BCB Insights for Poland",
    "Please illustrate the calculation of BCB in poland for different ranks",
    "I want to generate BCB requirements for Singapore after converting to Singapore currency. use the ranks from spain and bonus criteria and rates from Poland",
    "Create user stories for Poland Thermomix on the calculation of Basic Commission for advisors with EP and without EP, including examples. Specify the eligible ranks in both user stories.",
    "Create the user story for Poland Thermomix BCB the commission calculation in case of negative sales order event"
]

PCB_BCB_TEMPLATE = """
If the question is about providing / generating the requirements of a bonus, Please use the following section headers and content.
    Introduction
    - provide the bonus overview in this section
    Commission configurations
    - use the records with the type as "commissions configurations".
    - Show the parameter configuration and description as bullet points
    Rate card configurations
    - use the records with the type as "Rate card configurations".
    - Show the parameter configuration and description as bullet points
    Process
    - use the records with the type as "Process".
    - use the Key parameter definition text to describe the process

If the question is about Comparison of bonuses, show the output in Tabular format.

If the question is about generating requirements of a specific bonus or all the bonuses for a country, based on the parameters of existing countries, do the following:
- Convert the currency to the country for which the bonus requirements are being created
- Generate the bonus requirement as per the existing countries bonus structures.
"""

# Session state
if "messages" not in st.session_state:
    st.session_state.messages = []
if "uploaded_file" not in st.session_state:
    st.session_state.uploaded_file = None
if "extracted_text" not in st.session_state:
    st.session_state.extracted_text = ""
if "new_upload" not in st.session_state:
    st.session_state.new_upload = False
if "chat_sessions" not in st.session_state:
    st.session_state.chat_sessions = []
if "session_titles" not in st.session_state:
    st.session_state.session_titles = []
if "current_session_index" not in st.session_state:
    st.session_state.current_session_index = None

# File extraction
def extract_text(file):
    text = ""
    if file.type == "application/pdf":
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file.type == "text/csv":
        df = pd.read_csv(file)
        text += df.to_string()
    elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(file)
        text += df.to_string()
    return text.strip()

# Sidebar
with st.sidebar:
    st.title("How Can I Assist You? 🤖")

    if st.button("🆕 New Chat"):
        if st.session_state.messages:
            first_user_msg = next((msg["content"] for msg in st.session_state.messages if msg["role"] == "user"), "Untitled")
            st.session_state.chat_sessions.append(st.session_state.messages)
            st.session_state.session_titles.append(first_user_msg[:40] + "..." if len(first_user_msg) > 40 else first_user_msg)
        st.session_state.messages = []
        st.session_state.current_session_index = None
        st.rerun()

    if st.button("➕ Upload File"):
        st.session_state.new_upload = True

    if st.session_state.get("new_upload", False):
        uploaded_file = st.file_uploader("Upload File", type=["pdf", "docx", "pptx", "csv", "xlsx"], key="file_uploader")

        if uploaded_file:
            st.session_state.uploaded_file = uploaded_file.name
            st.session_state.extracted_text = extract_text(uploaded_file)
            st.session_state.new_upload = False

    if st.session_state.get("uploaded_file"):
        st.markdown(f"📂 **Uploaded File:** {st.session_state.uploaded_file}")
        if st.button("❌ Clear Uploaded File", use_container_width=True):
            st.session_state.uploaded_file = None
            st.session_state.extracted_text = ""

    if st.session_state.chat_sessions:
        st.markdown("---")
        st.subheader("🕓 Latest Chat")
        latest_index = len(st.session_state.chat_sessions) - 1
        title = st.session_state.session_titles[latest_index]
        if st.button(title, key=f"session_{latest_index}"):
            st.session_state.messages = st.session_state.chat_sessions[latest_index]
            st.session_state.current_session_index = latest_index
            st.rerun()

    # Renamed to "Prompts"
    st.markdown("---")
    st.subheader("Prompts")
    with st.expander("Choose a prompt and click to run it"):
        for idx, prompt in enumerate(PRESET_PROMPTS):
            if st.button(prompt, key=f"preset_{idx}"):
                st.session_state._selected_preset = prompt
                st.rerun()

# Display messages
for message in st.session_state.messages:
    if message["role"] != "system":
        with st.chat_message(message["role"]):
            if "file" in message and message["file"]:
                st.markdown(f"📄 **{message['file']}**")
            st.markdown(message["content"])

# OpenAI config
openai.api_key = "14560021aaf84772835d76246b53397a"
openai.api_base = 'https://amrxgenai.openai.azure.com/'
openai.api_type = 'azure'
openai.api_version = '2024-02-15-preview'
deployment_name = 'gpt'

# Process user input
def process_user_input(user_input_text):
    if not user_input_text or not user_input_text.strip():
        return

    chat_messages = st.session_state.messages.copy()
    pcb_bcb_in_query = "pcb" in user_input_text.lower() or "bcb" in user_input_text.lower()

    if st.session_state.extracted_text:
        if not any(m["role"] == "system" for m in chat_messages):
            system_prompt = f"Here is a document that provides context:\n\n{st.session_state.extracted_text}\n\n"
            if pcb_bcb_in_query:
                system_prompt += PCB_BCB_TEMPLATE + "\n\n"
            system_prompt += f"Now, based on this document, answer the following:\n{user_input_text}"
            chat_messages.insert(0, {"role": "system", "content": system_prompt})

    chat_messages.append({"role": "user", "content": user_input_text})

    try:
        response = openai.ChatCompletion.create(
            engine=deployment_name,
            messages=chat_messages,
            temperature=0,
            max_tokens=2000
        )
        ai_response = response["choices"][0]["message"]["content"]
    except Exception as e:
        ai_response = f"Error calling OpenAI: {e}"

    st.session_state.messages.append({
        "role": "user",
        "content": user_input_text,
        "file": st.session_state.uploaded_file
    })

    st.session_state.messages.append({
        "role": "assistant",
        "content": ai_response
    })

    if "_selected_preset" in st.session_state:
        del st.session_state["_selected_preset"]

    st.rerun()

# Chat input
user_input = st.chat_input("Enter your query:")
if user_input and user_input.strip():
    process_user_input(user_input)

if st.session_state.get("_selected_preset"):
    process_user_input(st.session_state["_selected_preset"])

# Download chat history
def create_docx():
    doc = Document()
    doc.add_heading("Chatbot Conversation History", level=1).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    for msg in st.session_state.messages:
        if msg["role"] == "system":
            continue
        p = doc.add_paragraph()
        run = p.add_run(msg["role"].capitalize() + "\n")
        run.bold = True
        run.font.size = Pt(14)
        if "file" in msg and msg["file"]:
            doc.add_paragraph(f"📄 {msg['file']}")
        p.add_run(msg["content"]).font.size = Pt(12)
        p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        doc.add_paragraph("\n")
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

if st.button("Download Chat History"):
    docx_file = create_docx()
    st.download_button(label="Download Chat History", data=docx_file, file_name="chat_history.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")