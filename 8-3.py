import streamlit as st
import openai
import docx
import pandas as pd
import PyPDF2
import pptx
from io import BytesIO
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

# OpenAI API Configuration
openai.api_key = "14560021aaf84772835d76246b53397a"
openai.api_base = "https://amrxgenai.openai.azure.com/"
openai.api_type = 'azure'
openai.api_version = '2024-02-15-preview'
deployment_name = 'gpt'

st.title("AI Chatbot with File Upload & Response History")

# Initialize session state for storing chat history
if "messages" not in st.session_state:
    st.session_state.messages = []

# Sidebar for user input & file upload
with st.sidebar:
    st.header("Chat Input & File Upload")
    user_input = st.text_input("Enter your query:", "", key="user_input")
    uploaded_file = st.file_uploader("Upload a file (PDF, DOCX, PPT, CSV, XLSX)", type=["pdf", "docx", "pptx", "csv", "xlsx"])
    send_button = st.button("Send")
    download_button = st.button("Download Chat History")

# Function to extract text from uploaded files
def extract_text(file):
    text = ""
    if file.type == "application/pdf":
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = pptx.Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # Check if shape has text
                    text += shape.text + "\n"
    elif file.type == "text/csv":
        df = pd.read_csv(file)
        text += df.to_string()
    elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(file)
        text += df.to_string()
    return text

# Process file upload if available
extracted_text = ""
if uploaded_file:
    extracted_text = extract_text(uploaded_file)
    st.text_area("Extracted Text:", extracted_text, height=150)

# Display chat history in main area
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if send_button and user_input:
    # Append user message to history
    prompt = user_input
    if extracted_text:
        prompt += "\n\n[Extracted Text]\n" + extracted_text
    
    st.session_state.messages.append({"role": "user", "content": user_input})
    
    # Generate response using OpenAI API
    response = openai.ChatCompletion.create(
        engine=deployment_name,
        messages=st.session_state.messages,
        temperature=0.7,
        max_tokens=2000
    )
    
    ai_response = response["choices"][0]["message"]["content"]
    
    # Append AI response to history
    st.session_state.messages.append({"role": "assistant", "content": ai_response})
    
    # Refresh page
    st.rerun()

# Function to create a DOCX file from chat history with the same format
def create_docx():
    doc = Document()
    doc.add_heading("Chatbot Conversation History", level=1).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    
    for msg in st.session_state.messages:
        p = doc.add_paragraph()
        run = p.add_run(msg["role"].capitalize() + "\n")
        run.bold = True
        run.font.size = Pt(14)
        p.add_run(msg["content"]).font.size = Pt(12)
        p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        doc.add_paragraph("\n")
    
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Download chat history as DOCX
if download_button:
    docx_file = create_docx()
    st.download_button(label="Download Chat History", data=docx_file, file_name="chat_history.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
